from pypdf import PdfReader
from docx import Document
from pptx import Presentation

def load_pdf(file_path):
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text

def load_docx(path):
    doc = Document(path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def load_pptx(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def load_txt(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()